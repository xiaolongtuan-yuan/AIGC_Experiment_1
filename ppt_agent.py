import os
import re
from string import Template

import openai
import json
import argparse
import time
import random
import datetime

from IPython.core.debugger import prompt
from certifi import contents
from pptx import Presentation
from pptx.util import Inches
from soupsieve import select

from draw_chart import csv_Line_chart, csv_pie_chart
from image_search import get_relative_images


# api_key = os.getenv("OPENAI_KEY")
# openai_endpoint = "https://api.aiproxy.io/v1"

class PPTAgent:
    def __init__(self, api_key, ollama_url=None, ollama_model=None):
        openai.api_key = api_key
        # openai.api_base = openai_endpoint

        self.ollama_url = ollama_url
        self.ollama_model = ollama_model

        self.tool_map = {
            "csv_Line_chart" : csv_Line_chart,
            "csv_pie_chart": csv_pie_chart
        }

        self.topic = ""
        self.file_paths = []
        self.charts = []

    def gen_ppt_content(self, topic, pages, language):
        self.topic = topic
        language_map = {"cn": "Chinese", "en": "English"}
        language_str = language_map[language]

        images = get_relative_images(topic)
        for image in images:
            self.charts.append(image)

        for file_path in self.file_paths:
            self.add_data_chart(file_path)

        output_format = self._get_output_format()
        messages = self._get_messages(topic, pages, language_str, output_format)
        print(messages)
        content = self._get_content(messages)
        return self._parse_content(content)

    def _get_output_format(self):
        output_format = open("resources/output_format.json", "r").read()
        return output_format

    def upload_files(self, file_path):
        self.file_paths.append(file_path)

    def add_data_chart(self, file_path):
        tool_system_prompt = open("resources/file_process_prompt.txt", "r").read()
        tool_system_prompt = Template(tool_system_prompt)
        with open("resources/draw_tools.json", "r") as tool_file:
            draw_tools = json.load(tool_file)

        prompt1 = tool_system_prompt.substitute(file_path=file_path, require=self.topic)
        message = [
            {
                "role": "user",
                "content": prompt1,
            }
        ]
        completion = openai.ChatCompletion.create(
            messages=message,
            model="gpt-3.5-turbo",
            tools=draw_tools,
            tool_choice='required'
        )
        func_call = completion.choices[0].message.tool_calls[0].function
        function_call = json.loads(str(func_call))
        func = self.tool_map[function_call['name']]
        res1 = func(**json.loads(function_call['arguments']))

        desciption_prompt = open("resources/file_describe_prompt.txt", "r").read()
        desciption_prompt = Template(desciption_prompt)
        data_str = open(file_path, "r").read()
        prompt2 = desciption_prompt.substitute(file_path=file_path, data=data_str)
        completion2 = openai.ChatCompletion.create(
            messages=[{
                "role": "user",
                "content": prompt2,
            }],
            model="gpt-3.5-turbo",
        )
        response = completion2.choices[0].message.content

        self.charts.append({
            "file_path":res1,
            "data_insights":response
        })
        print("process file successfully")


    def _get_messages(self, topic, pages, language_str, output_format):
        content = open("resources/ppt_content_prompt.txt", "r").read()
        content = Template(content)
        chart_str = str(self.charts)
        prompt = content.substitute(topic=topic, output_format=output_format, chart=chart_str, pages=pages, language_str=language_str)
        return [
            {
                "role": "user",
                "content": prompt,
            }
        ]

    def _get_content(self, messages):
        completion = openai.ChatCompletion.create(
            model="gpt-3.5-turbo", messages=messages
        )
        return completion.choices[0].message.content

    def _parse_content(self, content):
        try:
            match = re.search(r"(\{.*\})", content, re.DOTALL)
            if match:
                content = match.groups()[0]
            return json.loads(content.strip())
        except Exception as e:
            print(f"The response is not a valid JSON format: {e}")
            print("I'm a PPT assistant, your PPT generate failed, please retry later..")
            raise Exception("The LLM return invalid result, please retry later..")
            exit(1)

    def generate_ppt(self, content, template=None):
        ppt = Presentation()
        if template:
            ppt = Presentation(template)
        self.create_slides(ppt, content)
        ppt_name = self._get_ppt_name(content)
        save_path = os.path.join("files", ppt_name)
        ppt.save(save_path)
        return ppt_name

    def create_slides(self, presentation, content):
        self.create_title_slide(presentation, content)
        pages = content.get("pages", [])
        for index, page in enumerate(pages):
            self.create_content_slide(presentation, page, index)

    def create_title_slide(self, presentation, content):
        title_slide_layout = presentation.slide_layouts[0]
        title_slide = presentation.slides.add_slide(title_slide_layout)
        title_slide.shapes.title.text = content.get("title", "")
        title_slide.placeholders[1].text = "Generated by PPT Agent"

    def create_content_slide(self, presentation, page, index):
        page_type = page.get("type")
        if page.get("type") == "chart":
            page_title = page.get("title", "")

            # 创建使用标题和内容的幻灯片布局
            chart_slide_layout = presentation.slide_layouts[1]
            chart_slide = presentation.slides.add_slide(chart_slide_layout)

            # 设置幻灯片标题
            chart_slide.shapes.title.text = page_title

            # 添加数据描述
            content = page.get("content", [])
            if content:
                # 获取描述和文件路径
                data_description = content[0].get("description", "")
                chart_path = content[0].get("file_path", "")

                # 添加数据描述文本框
                textbox = chart_slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(1))
                text_frame = textbox.text_frame
                text_frame.text = data_description

                # 如果文件路径有效，则添加图片
                if chart_path:
                    left = Inches(5)  # 设置图片位置
                    top = Inches(2)
                    chart_slide.shapes.add_picture(chart_path, left, top, width=Inches(4))

        else:
            page_title = page.get("title", "")
            bullet_slide_layout = presentation.slide_layouts[1]
            bullet_slide = presentation.slides.add_slide(bullet_slide_layout)
            bullet_slide.shapes.title.text = page_title
            self.add_bullets_to_slide(bullet_slide, page)

    def add_bullets_to_slide(self, slide, page):
        body_shape = slide.shapes.placeholders[1]
        for bullet in page.get("content", []):
            self.add_bullet(body_shape, bullet)

    def add_bullet(self, body_shape, bullet):
        bullet_title = bullet.get("title", "")
        bullet_description = bullet.get("description", "")
        self.add_paragraph(body_shape, bullet_title, level=1)
        self.add_paragraph(body_shape, bullet_description, level=2)

    def add_paragraph(self, body_shape, text, level):
        paragraph = body_shape.text_frame.add_paragraph()
        paragraph.text = text
        paragraph.level = level

    def _get_ppt_name(self, content):
        ppt_name = content.get("title", "")
        ppt_name = re.sub(r'[\\/:*?"<>|]', "", ppt_name)
        timestamp = datetime.datetime.now().strftime("%Y%m%d%H%M%S")
        return f"{ppt_name}_{timestamp}.pptx"
